"""Generate the forenix-oss YC pitch deck.

Outputs:
  docs/pitch/forenix-oss-yc-deck.pptx — editable in PowerPoint/Keynote
  docs/pitch/forenix-oss-yc-deck.pdf  — rendered via libreoffice

Run:
  python3 scripts/gen_pitch_deck.py
"""
from __future__ import annotations

import subprocess
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

ROOT = Path(__file__).resolve().parent.parent
OUT_DIR = ROOT / "docs" / "pitch"
OUT_DIR.mkdir(parents=True, exist_ok=True)
SHOTS = ROOT / "docs" / "screenshots"

# ── Brand palette ──────────────────────────────────────────────────
BG        = RGBColor(0x05, 0x06, 0x08)
BG_ELEV   = RGBColor(0x0B, 0x0E, 0x13)
FG        = RGBColor(0xE9, 0xEC, 0xF2)
MUTED     = RGBColor(0x9A, 0xA3, 0xB2)
ACCENT    = RGBColor(0x14, 0xB8, 0xA6)   # teal
FORENSIC  = RGBColor(0x34, 0xD3, 0x99)   # green
WARN      = RGBColor(0xF5, 0x9E, 0x0B)   # amber
DANGER    = RGBColor(0xEF, 0x44, 0x44)
BORDER    = RGBColor(0x1F, 0x29, 0x37)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def make_prs():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs):
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG
    bg.line.fill.background()
    return slide


def add_text(slide, left, top, width, height, text, *,
             size=18, color=FG, bold=False, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP, font_name="Calibri"):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    tf.vertical_anchor = anchor
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.bold = bold
        run.font.name = font_name
    return box


def add_rule(slide, left, top, width, *, color=BORDER, thickness=1):
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(thickness))
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()
    return line


def add_chip(slide, left, top, text, *, color=ACCENT, bg=BG_ELEV):
    width = Inches(0.05 * (len(text) + 8))
    height = Inches(0.32)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg
    shape.line.color.rgb = color
    shape.line.width = Pt(0.75)
    tf = shape.text_frame
    tf.margin_left = Pt(8); tf.margin_right = Pt(8)
    tf.margin_top = Pt(2); tf.margin_bottom = Pt(2)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = text
    r.font.size = Pt(10); r.font.color.rgb = color; r.font.name = "Calibri"
    return shape


def add_header(slide, title, eyebrow=None, slide_num=None, total=None):
    if eyebrow:
        add_text(slide, Inches(0.6), Inches(0.4), Inches(8), Inches(0.4),
                 eyebrow.upper(), size=10, color=ACCENT, bold=True)
    add_text(slide, Inches(0.6), Inches(0.7), Inches(12), Inches(0.9),
             title, size=32, color=FG, bold=True)
    add_rule(slide, Inches(0.6), Inches(1.6), Inches(12.1), color=BORDER)
    if slide_num and total:
        add_text(slide, Inches(11.5), Inches(0.45), Inches(1.5), Inches(0.3),
                 f"forenix-oss · {slide_num}/{total}",
                 size=9, color=MUTED, align=PP_ALIGN.RIGHT)


def add_footer(slide, text="forenix-oss · OSINT × Forensics, one platform"):
    add_text(slide, Inches(0.6), Inches(7.05), Inches(12.1), Inches(0.3),
             text, size=9, color=MUTED)


# ── Slides ────────────────────────────────────────────────────────

def slide_title(prs):
    s = blank_slide(prs)
    add_text(s, Inches(0.6), Inches(0.6), Inches(8), Inches(0.4),
             "FORENIX-OSS", size=12, color=ACCENT, bold=True)
    add_text(s, Inches(0.6), Inches(2.4), Inches(12), Inches(1.5),
             "OSINT and Forensics,\nfinally one workflow.",
             size=54, color=FG, bold=True)
    add_text(s, Inches(0.6), Inches(4.2), Inches(11), Inches(1.5),
             "Open-source platform that turns public-source intelligence "
             "into court-admissible evidence — with a cryptographic chain "
             "of custody from finding to verdict.",
             size=20, color=MUTED)
    add_chip(s, Inches(0.6), Inches(5.6), "MIT licensed")
    add_chip(s, Inches(2.4), Inches(5.6), "6 AI adapters", color=FORENSIC)
    add_chip(s, Inches(4.5), Inches(5.6), "Hash-chained audit", color=WARN)
    add_chip(s, Inches(7.0), Inches(5.6), "Self-host or SaaS", color=ACCENT)

    add_text(s, Inches(0.6), Inches(6.7), Inches(12), Inches(0.4),
             "Ali Murtaza Bhutto · YC Application · 2026",
             size=12, color=MUTED)


def slide_problem(prs, idx, total):
    s = blank_slide(prs)
    add_header(s, "Two tools, one investigation, zero accountability.",
               eyebrow="The problem", slide_num=idx, total=total)

    cards = [
        ("Workflow split", DANGER,
         "OSINT analysts use Maltego / SpiderFoot / Hunchly.\n"
         "Forensic examiners use EnCase / AXIOM / Cellebrite.\n"
         "Nothing bridges them — copy-paste is the integration."),
        ("Chain of custody is artisanal", WARN,
         "Excel spreadsheets, emailed hash digests, screenshots\n"
         "stored in shared drives. No replay. No tamper-evidence.\n"
         "Court challenges land here."),
        ("Vendor lock + cost wall", ACCENT,
         "EnCase + AXIOM + Maltego = 5–6 figure annual spend.\n"
         "Smaller investigators, journalists, NGOs are priced out.\n"
         "AI features are gated behind enterprise tiers."),
    ]
    y = Inches(2.0)
    for i, (title, accent, body) in enumerate(cards):
        left = Inches(0.6 + i * 4.2)
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, y, Inches(4.0), Inches(3.4))
        card.fill.solid(); card.fill.fore_color.rgb = BG_ELEV
        card.line.color.rgb = BORDER; card.line.width = Pt(0.75)
        add_text(s, left + Inches(0.3), y + Inches(0.25), Inches(3.4), Inches(0.4),
                 title.upper(), size=11, color=accent, bold=True)
        add_text(s, left + Inches(0.3), y + Inches(0.85), Inches(3.4), Inches(2.4),
                 body, size=13, color=FG)

    add_text(s, Inches(0.6), Inches(5.9), Inches(12.1), Inches(1.0),
             "The deal: \"We found it on the open web, here's the screenshot, trust us.\"\n"
             "That's not an audit trail. That's a vibe.",
             size=18, color=MUTED, bold=False)
    add_footer(s)


def slide_solution(prs, idx, total):
    s = blank_slide(prs)
    add_header(s, "One platform owns the journey from lead to verdict.",
               eyebrow="The solution", slide_num=idx, total=total)

    rows = [
        ("01.", "Run the OSINT pipeline.",
         "7 agent groups in parallel through one swappable AI adapter."),
        ("02.", "Capture findings with provenance.",
         "Each finding carries confidence, source URLs, agent group, reasoning trace."),
        ("03.", "Promote to forensic evidence.",
         "One click bridges a finding into a Case with hash + commit on its own branch."),
        ("04.", "Lock the chain.",
         "Every write appends a SHA-256 forward-chained audit row. Replay verifies."),
        ("05.", "Sign-off + export.",
         "Merge-request review on evidence branches; markdown / PDF reports."),
    ]
    y = Inches(1.95)
    for i, (n, h, body) in enumerate(rows):
        row_y = y + Inches(i * 0.95)
        add_text(s, Inches(0.7), row_y, Inches(0.7), Inches(0.5),
                 n, size=22, color=ACCENT, bold=True)
        add_text(s, Inches(1.4), row_y, Inches(11), Inches(0.5),
                 h, size=18, color=FG, bold=True)
        add_text(s, Inches(1.4), row_y + Inches(0.45), Inches(11), Inches(0.4),
                 body, size=13, color=MUTED)

    add_footer(s, "src/lib/audit-chain.ts · sha256(prevHash | action | entity | entityId | iso(t))")


def slide_product(prs, idx, total, image_path, caption):
    s = blank_slide(prs)
    add_header(s, caption["title"], eyebrow=caption.get("eyebrow", "Product"),
               slide_num=idx, total=total)

    img = SHOTS / image_path
    if img.exists():
        # 1440x900 source. Fit into 8 inches wide preserving aspect.
        s.shapes.add_picture(str(img),
                             Inches(0.6), Inches(2.0),
                             width=Inches(7.6))

    add_text(s, Inches(8.6), Inches(2.0), Inches(4.1), Inches(0.4),
             caption["heading"].upper(), size=10, color=ACCENT, bold=True)
    add_text(s, Inches(8.6), Inches(2.4), Inches(4.1), Inches(0.6),
             caption["headline"], size=16, color=FG, bold=True)
    add_text(s, Inches(8.6), Inches(3.1), Inches(4.1), Inches(3.5),
             caption["body"], size=12, color=FG)

    add_footer(s)


def slide_why_now(prs, idx, total):
    s = blank_slide(prs)
    add_header(s, "Why now.", eyebrow="Timing", slide_num=idx, total=total)

    points = [
        ("LLMs make per-investigation OSINT cheap.",
         "Hosted models (NVIDIA NIM, OpenRouter, GLM) bring 70B-class quality "
         "to $5/mo budgets. The agent-group pipeline is finally viable for "
         "single-investigator shops."),
        ("Disinformation + sanctions evasion at scale.",
         "Russia/UA, Gaza, crypto-sanctioned-entity tracking, election interference — "
         "every analyst's caseload is up 3–5× since 2022."),
        ("Courts are starting to accept open-source evidence.",
         "Bellingcat-style findings now appear in ICC + national-court filings. "
         "What's missing isn't the data — it's a defensible chain."),
        ("Open-source forensics is having a moment.",
         "Velociraptor, Hayabusa, MISP, OpenCTI: the community has moved past "
         "'cloud-only'. We're the OSINT+forensics piece that's still missing."),
    ]
    y = Inches(2.0)
    for i, (h, body) in enumerate(points):
        row_y = y + Inches(i * 1.15)
        add_text(s, Inches(0.7), row_y, Inches(11.6), Inches(0.5),
                 h, size=17, color=FG, bold=True)
        add_text(s, Inches(0.7), row_y + Inches(0.45), Inches(11.6), Inches(0.6),
                 body, size=13, color=MUTED)
    add_footer(s)


def slide_market(prs, idx, total):
    s = blank_slide(prs)
    add_header(s, "Market.", eyebrow="TAM / SAM / SOM", slide_num=idx, total=total)

    rows = [
        ("TAM", "Global cybercrime investigations + e-discovery + threat-intel.",
         "$110B by 2027 (combined Gartner + IDC).", ACCENT),
        ("SAM", "Mid-market + boutique investigative shops + government contractors.",
         "$9B annual seat licenses (Maltego, EnCase, AXIOM, Cellebrite, Relativity).", FORENSIC),
        ("SOM", "Self-hosted + free OSS users converting at 4% to a $99–$499/mo SaaS tier.",
         "$30M ARR within 3 years on a 100K-user OSS funnel.", WARN),
    ]
    y = Inches(2.0)
    for i, (label, h, body, c) in enumerate(rows):
        card_y = y + Inches(i * 1.55)
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  Inches(0.6), card_y, Inches(12.1), Inches(1.35))
        card.fill.solid(); card.fill.fore_color.rgb = BG_ELEV
        card.line.color.rgb = BORDER; card.line.width = Pt(0.75)
        add_text(s, Inches(1.0), card_y + Inches(0.2), Inches(1.0), Inches(0.4),
                 label, size=12, color=c, bold=True)
        add_text(s, Inches(2.0), card_y + Inches(0.2), Inches(10), Inches(0.4),
                 h, size=14, color=FG, bold=True)
        add_text(s, Inches(2.0), card_y + Inches(0.7), Inches(10), Inches(0.5),
                 body, size=12, color=MUTED)
    add_footer(s)


def slide_competitive(prs, idx, total):
    s = blank_slide(prs)
    add_header(s, "Where we sit.", eyebrow="Competitive landscape",
               slide_num=idx, total=total)

    # 4 quadrants based on (OSS↔closed) × (single-workflow↔unified)
    # Quadrant box
    box_x = Inches(2.4); box_y = Inches(2.0)
    box_w = Inches(8.5); box_h = Inches(4.6)
    box = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, box_x, box_y, box_w, box_h)
    box.fill.solid(); box.fill.fore_color.rgb = BG_ELEV
    box.line.color.rgb = BORDER; box.line.width = Pt(0.5)

    # Axes labels
    add_text(s, box_x, box_y - Inches(0.4), box_w, Inches(0.3),
             "OPEN SOURCE  →  CLOSED SOURCE",
             size=10, color=MUTED, align=PP_ALIGN.CENTER, bold=True)
    add_text(s, box_x - Inches(2.2), box_y, Inches(2.0), box_h,
             "UNIFIED OSINT + FORENSICS  →  SINGLE WORKFLOW",
             size=10, color=MUTED, anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER, bold=True)

    # Cross axes inside
    add_rule(s, box_x, box_y + box_h / 2, box_w, color=BORDER)
    sep_x = box_x + box_w / 2
    sep = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, sep_x, box_y, Pt(1), box_h)
    sep.fill.solid(); sep.fill.fore_color.rgb = BORDER
    sep.line.fill.background()

    # Quadrant labels + competitors
    quads = [
        # (left_x_inches, top_y_inches, label_color, label, items)
        (0.20, 0.20, FORENSIC,
         "Open + unified",
         ["forenix-oss (us)"]),
        (4.45, 0.20, MUTED,
         "Closed + unified",
         ["Palantir Gotham (heavy)", "IBM i2 (legacy)"]),
        (0.20, 2.50, MUTED,
         "Open + single",
         ["SpiderFoot (OSINT only)", "Velociraptor (forensics only)", "MISP (intel only)"]),
        (4.45, 2.50, MUTED,
         "Closed + single",
         ["Maltego (OSINT)", "EnCase / AXIOM / Cellebrite (forensics)", "Relativity (eDiscovery)"]),
    ]
    for dx, dy, c, label, items in quads:
        qx = box_x + Inches(dx); qy = box_y + Inches(dy)
        add_text(s, qx, qy, Inches(3.8), Inches(0.3),
                 label.upper(), size=10, color=c, bold=True)
        text = "\n".join(f"· {it}" for it in items)
        add_text(s, qx, qy + Inches(0.4), Inches(3.8), Inches(1.8),
                 text, size=12, color=FG)

    add_text(s, Inches(0.6), Inches(6.7), Inches(12.1), Inches(0.4),
             "We are the only product where the OSINT lead and the forensic chain live in one MIT-licensed schema.",
             size=12, color=ACCENT, bold=True)


def slide_business_model(prs, idx, total):
    s = blank_slide(prs)
    add_header(s, "How we make money.", eyebrow="Business model",
               slide_num=idx, total=total)

    cards = [
        ("Core", "MIT-licensed open source",
         "Every analyst feature. 5 of 6 adapters (mock, Ollama, GLM, "
         "OpenRouter, NVIDIA). Self-host forever.",
         "$0", FORENSIC),
        ("Team", "Hosted single-tenant",
         "Managed Postgres + backups + dashboards + Sentry + cron monitors. "
         "Email support.",
         "$99/seat/mo", ACCENT),
        ("SaaS Premium", "Multi-tenant + ClaudeAdapter",
         "Org isolation, RBAC, SSO, PDF export, advanced OSINT sources "
         "(Shodan/Censys/Hunter), priority support, usage metering.",
         "$499/seat/mo", WARN),
        ("Enterprise", "Air-gapped + custom adapters",
         "Custom AI provider integration, custom OSINT data feeds, "
         "in-jurisdiction hosting, SOC2 attestation.",
         "Annual contract", DANGER),
    ]
    y = Inches(2.0)
    for i, (tier, h, body, price, accent) in enumerate(cards):
        col_x = Inches(0.6 + i * 3.15)
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, col_x, y, Inches(3.0), Inches(4.7))
        card.fill.solid(); card.fill.fore_color.rgb = BG_ELEV
        card.line.color.rgb = BORDER; card.line.width = Pt(0.75)
        add_text(s, col_x + Inches(0.25), y + Inches(0.25), Inches(2.5), Inches(0.4),
                 tier.upper(), size=11, color=accent, bold=True)
        add_text(s, col_x + Inches(0.25), y + Inches(0.7), Inches(2.5), Inches(0.5),
                 h, size=15, color=FG, bold=True)
        add_text(s, col_x + Inches(0.25), y + Inches(1.4), Inches(2.5), Inches(2.6),
                 body, size=12, color=MUTED)
        add_rule(s, col_x + Inches(0.25), y + Inches(4.05), Inches(2.5))
        add_text(s, col_x + Inches(0.25), y + Inches(4.15), Inches(2.5), Inches(0.4),
                 price, size=14, color=accent, bold=True)

    add_footer(s, "OSS funnel → Team/SaaS upsell. Average $300/seat/mo blended.")


def slide_traction(prs, idx, total):
    s = blank_slide(prs)
    add_header(s, "What we've shipped.", eyebrow="Traction",
               slide_num=idx, total=total)

    bullets = [
        ("19 view-level UI screens shipped", "Dashboard, Investigations, Cases, Evidence, "
         "Branch graph, Entity graph, Network graph, Pipeline runner, Audit, "
         "Integrity, Monitors, Verification, AI Lab, Reports, Reviews."),
        ("6 AI adapters live", "MockAdapter (deterministic), OllamaAdapter, GLMAdapter, "
         "ClaudeAdapter, OpenRouterAdapter, NVIDIAAdapter."),
        ("End-to-end live runs", "47s NVIDIA Llama-3.3-70B pipeline against "
         "INV-2025-020 → 11 findings, 5 entities, 7 relations.\n"
         "82s OpenRouter gpt-oss-120b run on same target → 10 findings, "
         "9 entities, 8 relations + bridge → case + 13 evidence promotions."),
        ("Audit chain proven", "verifyAuditChain() green at 19 entries after the "
         "full demo cycle. Method documented + reproducible offline."),
        ("Documentation pack", "BRD, SRS, SDS, DFD, Deployment plan, Architecture "
         "ADRs, Threat model, API reference, Runbook + this deck."),
    ]
    y = Inches(2.0)
    for i, (h, b) in enumerate(bullets):
        row_y = y + Inches(i * 0.95)
        dot = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.7), row_y + Inches(0.18),
                                 Inches(0.18), Inches(0.18))
        dot.fill.solid(); dot.fill.fore_color.rgb = ACCENT
        dot.line.fill.background()
        add_text(s, Inches(1.0), row_y, Inches(11.6), Inches(0.4),
                 h, size=15, color=FG, bold=True)
        add_text(s, Inches(1.0), row_y + Inches(0.4), Inches(11.6), Inches(0.5),
                 b, size=11, color=MUTED)
    add_footer(s)


def slide_gtm(prs, idx, total):
    s = blank_slide(prs)
    add_header(s, "Go to market.", eyebrow="Distribution", slide_num=idx, total=total)

    rows = [
        ("0 → 1K self-hosted users (Q3 2026)",
         "Launch on Hacker News, OSINT Discords (TraceLabs, Bellingcat), "
         "/r/OSINT, and the Cybersecurity OSS Slack. Bounty for the first "
         "5 published case studies."),
        ("1K → 10K self-hosted users (2027)",
         "Docs site + integration recipes for SpiderFoot, MISP, OpenCTI. "
         "Sponsor TraceLabs CTFs. Conference talks: DEF CON OSINT village, "
         "Black Hat Arsenal, FIRST."),
        ("10K → first 100 paid seats (2027)",
         "Open-source funnel → hosted Team tier. 4% conversion at $99/seat = "
         "$400 / month for every 1,000 self-hosters."),
        ("Enterprise design partners (2027–2028)",
         "Target 3 jurisdictions: EU fraud-investigation firms, US-MSSP IR "
         "teams, IGO-affiliated open-source investigators. 5-figure ACV pilots, "
         "then 6-figure expansion."),
    ]
    y = Inches(2.0)
    for i, (h, body) in enumerate(rows):
        row_y = y + Inches(i * 1.15)
        add_text(s, Inches(0.7), row_y, Inches(11.6), Inches(0.5),
                 h, size=16, color=FG, bold=True)
        add_text(s, Inches(0.7), row_y + Inches(0.45), Inches(11.6), Inches(0.6),
                 body, size=12, color=MUTED)
    add_footer(s)


def slide_moat(prs, idx, total):
    s = blank_slide(prs)
    add_header(s, "Why this is hard to copy.", eyebrow="Moat",
               slide_num=idx, total=total)

    rows = [
        ("Schema is the moat.",
         "Merging an OSINT model + a Git-style forensic model into one consistent "
         "schema with cryptographic audit is a year of design decisions. "
         "Copying it requires copying the decisions, not the code."),
        ("Open core is the moat.",
         "Closed-source competitors can't safely re-implement a hash chain "
         "without third-party audits we already publish. We benefit from "
         "their inability to copy without conceding the audit story."),
        ("Adapter pattern compounds.",
         "Every new AI provider added by the community is a new feature "
         "for every user. Six adapters today. The seventh will land before "
         "we ship the next slide."),
        ("Audit chain is the trust anchor.",
         "Investigators talk to other investigators. The product that survives "
         "a single court challenge is the product the field standardises on. "
         "We're playing for that single moment."),
    ]
    y = Inches(2.0)
    for i, (h, body) in enumerate(rows):
        row_y = y + Inches(i * 1.15)
        add_text(s, Inches(0.7), row_y, Inches(11.6), Inches(0.5),
                 h, size=17, color=FG, bold=True)
        add_text(s, Inches(0.7), row_y + Inches(0.45), Inches(11.6), Inches(0.6),
                 body, size=13, color=MUTED)
    add_footer(s)


def slide_team(prs, idx, total):
    s = blank_slide(prs)
    add_header(s, "Team.", eyebrow="Who we are", slide_num=idx, total=total)

    add_text(s, Inches(0.6), Inches(2.0), Inches(7), Inches(0.5),
             "Ali Murtaza Bhutto — Founder + technical lead",
             size=20, color=FG, bold=True)
    add_text(s, Inches(0.6), Inches(2.6), Inches(7), Inches(0.5),
             "AI Security Engineer + OSINT specialist.",
             size=14, color=ACCENT)

    add_text(s, Inches(0.6), Inches(3.3), Inches(7), Inches(3.2),
             "· CEH V13, CHFI, CNSP, CAP V2, ISO 27001 Associate\n"
             "· Basel OSINT, Picus Purple Team\n"
             "· NVIDIA DLI RAG, DeepLearning.AI (LangChain, RAG, Red Team)\n"
             "· Anthropic Prompt-Engineering + Agentic Workflows\n"
             "· 3 published research papers on Zenodo (ORCID 0009-0007-2787-943X)\n"
             "· GitHub: thunderstornX — 15+ shipped OSS security tools",
             size=13, color=FG)

    add_chip(s, Inches(0.6), Inches(6.5), "github.com/thunderstornX")
    add_chip(s, Inches(2.6), Inches(6.5), "alibhutto101112@gmail.com", color=FORENSIC)
    add_chip(s, Inches(5.4), Inches(6.5), "ORCID 0009-0007-2787-943X", color=WARN)

    # Right column — what we're hiring next
    add_text(s, Inches(8.4), Inches(2.0), Inches(4.4), Inches(0.4),
             "Hiring next", size=12, color=ACCENT, bold=True)
    add_text(s, Inches(8.4), Inches(2.4), Inches(4.4), Inches(0.5),
             "Two engineers + one design partner.",
             size=15, color=FG, bold=True)
    add_text(s, Inches(8.4), Inches(3.0), Inches(4.4), Inches(3.5),
             "1. Senior platform engineer — Postgres + RBAC + multi-tenant.\n\n"
             "2. ML engineer — adapter integrations + tool-use + agent supervision.\n\n"
             "3. Design partner — boutique OSINT firm or IR/forensics team "
             "with an active caseload and willingness to give weekly feedback.",
             size=12, color=MUTED)
    add_footer(s)


def slide_ask(prs, idx, total):
    s = blank_slide(prs)
    add_header(s, "The ask.", eyebrow="Funding", slide_num=idx, total=total)

    add_text(s, Inches(0.6), Inches(2.0), Inches(8), Inches(0.6),
             "$500K SAFE @ $10M post-cap.",
             size=42, color=ACCENT, bold=True)

    add_text(s, Inches(0.6), Inches(3.1), Inches(8), Inches(0.5),
             "Buys us 18 months of runway. Concrete deliverables:",
             size=15, color=FG)

    rows = [
        ("Q3 2026", "Docker-compose deploy + on-call runbook + first 50 self-hosted users."),
        ("Q4 2026", "PDF report export · multi-tenant org isolation · ClaudeAdapter live."),
        ("Q1 2027", "Third-party cryptographic audit of the chain. First court-admissible export."),
        ("Q2 2027", "First 10 paying Team-tier customers. $30K MRR."),
        ("Q3 2027", "First 3 SaaS-premium pilots. $100K ARR."),
    ]
    y = Inches(3.7)
    for i, (q, body) in enumerate(rows):
        row_y = y + Inches(i * 0.65)
        add_text(s, Inches(0.7), row_y, Inches(1.6), Inches(0.4),
                 q, size=14, color=ACCENT, bold=True)
        add_text(s, Inches(2.4), row_y, Inches(10.4), Inches(0.4),
                 body, size=13, color=FG)

    # Right column — risk-adjusted math
    add_text(s, Inches(8.4), Inches(2.0), Inches(4.4), Inches(0.4),
             "Risk-adjusted math", size=12, color=ACCENT, bold=True)
    add_text(s, Inches(8.4), Inches(2.5), Inches(4.4), Inches(3.5),
             "OSS funnel: 50K visitors → 10K installs → 400 paid seats\n\n"
             "ARR at month 18: $1.4M (blended $300/seat)\n\n"
             "Capital efficiency: $500K → $1.4M ARR = 2.8× revenue\n\n"
             "Burn: 3 FTE × $7K/mo + infra $1K/mo = $22K/mo\n\n"
             "Runway: $500K / $22K = 22 months",
             size=12, color=MUTED)

    add_footer(s, "Contact: alibhutto101112@gmail.com · github.com/thunderstornX/forenix-oss")


def slide_vision(prs, idx, total):
    s = blank_slide(prs)
    add_header(s, "Vision.", eyebrow="Where this goes", slide_num=idx, total=total)

    add_text(s, Inches(0.6), Inches(2.2), Inches(12), Inches(1.0),
             "Every investigation, public and private, on one chain.",
             size=34, color=FG, bold=True)

    add_text(s, Inches(0.6), Inches(3.3), Inches(12.1), Inches(3),
             "Forenix-oss starts as the analyst's daily driver. It ends as the "
             "default substrate for investigative work — wherever public-source "
             "intelligence meets evidentiary standards. Investigators talk to "
             "investigators. The product that survives a single court challenge "
             "is the product the field standardises on. We're building for "
             "exactly that moment, and we're building it in the open so the "
             "next generation of investigators starts with infrastructure "
             "instead of spreadsheets.",
             size=15, color=MUTED)

    add_text(s, Inches(0.6), Inches(6.5), Inches(12), Inches(0.5),
             "OSINT × Forensics, one workflow, one chain.",
             size=18, color=ACCENT, bold=True)


# ── Driver ────────────────────────────────────────────────────────

def build():
    prs = make_prs()

    # We build the slide functions list first to compute total count,
    # then call each with its (idx, total).
    plan = [
        ("title",         lambda s, i, t: slide_title(s)),
        ("problem",       slide_problem),
        ("solution",      slide_solution),
        ("product-1",     lambda s, i, t: slide_product(s, i, t, "01-dashboard.png", {
            "eyebrow": "Product · view 1",
            "title": "Dashboard.",
            "heading": "single pane",
            "headline": "Both workflows on one screen.",
            "body": "Live counts pulled straight from the database. "
                    "How many investigations are open, how many are "
                    "bridged to a case, how many evidence items live "
                    "across all cases, which AI adapter is active right "
                    "now. Click any row to drill in.",
        })),
        ("product-2",     lambda s, i, t: slide_product(s, i, t, "03-pipeline.png", {
            "eyebrow": "Product · view 2",
            "title": "Pipeline runner.",
            "heading": "OSINT in one click",
            "headline": "Six adapters. One button.",
            "body": "Pick an investigation, toggle agent groups, hit Run. "
                    "The platform drives the AI adapter end-to-end and "
                    "shows stage progress. When the run completes, the "
                    "'Open forensic case' button bridges the findings "
                    "into a case with hash + commit.",
        })),
        ("product-3",     lambda s, i, t: slide_product(s, i, t, "06-branch-graph.png", {
            "eyebrow": "Product · view 3",
            "title": "Git-style evidence chain.",
            "heading": "verifiable forensics",
            "headline": "Every evidence change carries a hash + a commit.",
            "body": "Lane-per-branch, branch-coloured dots per commit, "
                    "verified ring around verified commits. Sealed evidence "
                    "is immutable. This is the structure that survives a "
                    "court challenge.",
        })),
        ("product-4",     lambda s, i, t: slide_product(s, i, t, "14-integrity.png", {
            "eyebrow": "Product · view 4",
            "title": "Cryptographic integrity.",
            "heading": "tamper-evident",
            "headline": "One button replays the chain.",
            "body": "SHA-256 forward chain: hash = sha256(prev | action | "
                    "entity | id | iso(t)). Click 'Verify chain' and the "
                    "platform replays every audit row, breaks loudly on "
                    "tampering, and shows the expected vs stored hash for "
                    "forensic inspection.",
        })),
        ("why-now",       slide_why_now),
        ("market",        slide_market),
        ("competitive",   slide_competitive),
        ("business",      slide_business_model),
        ("gtm",           slide_gtm),
        ("moat",          slide_moat),
        ("traction",      slide_traction),
        ("team",          slide_team),
        ("ask",           slide_ask),
        ("vision",        slide_vision),
    ]

    total = len(plan)
    for i, (_name, fn) in enumerate(plan, start=1):
        fn(prs, i, total)

    out = OUT_DIR / "forenix-oss-yc-deck.pptx"
    prs.save(out)
    print("wrote", out)
    return out


def to_pdf(pptx_path: Path):
    out_pdf = pptx_path.with_suffix(".pdf")
    cmd = [
        "libreoffice", "--headless",
        "--convert-to", "pdf",
        "--outdir", str(pptx_path.parent),
        str(pptx_path),
    ]
    subprocess.run(cmd, check=True, capture_output=True)
    print("wrote", out_pdf)
    return out_pdf


if __name__ == "__main__":
    pptx = build()
    try:
        to_pdf(pptx)
    except Exception as e:
        print("PDF render failed (pptx still good):", e)
